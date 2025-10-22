import os
from pathlib import Path
import sys
BASE = Path(__file__).parent.resolve()
PYTHON_EXE = sys.executable
EXTRACT = (BASE / "extraction.py").resolve()

assert Path(PYTHON_EXE).exists(), f"Bad PYTHON_EXE: {PYTHON_EXE}"
assert EXTRACT.exists(), f"Missing script: {EXTRACT}"

import argparse
from typing import List

from pdfminer.high_level import extract_text
import re
from pptx import Presentation
from collections import Counter
import unicodedata

def sanitize(s: str) -> str:
    # Normalize Unicode (e.g., fancy bullets/quotes)
    s = unicodedata.normalize("NFKC", s)

    # Standardize newlines
    s = s.replace("\r\n", "\n").replace("\r", "\n")

    # Replace vertical tab & form feed with newlines
    s = s.replace("\x0b", "\n").replace("\x0c", "\n")

    # Remove all other C0 control chars except \n and \t
    s = re.sub(r"[\x00-\x08\x0E-\x1F\x7F]", "", s)

    # Remove zero-width and BOM characters
    s = s.replace("\u200b", "")  # zero-width space
    s = s.replace("\u200c", "").replace("\u200d", "")  # ZWNJ/ZWJ
    s = s.replace("\ufeff", "")  # BOM

    # Collapse >2 consecutive blank lines
    s = re.sub(r"\n{3,}", "\n\n", s)

    # Trim trailing spaces on lines
    s = "\n".join(line.rstrip() for line in s.splitlines())

    return s

def drop_headers_footers(text: str, repeat_threshold: int = 3) -> str:
    """
    Heuristic: any short line that repeats >= repeat_threshold times is likely a header/footer.
    Also drops slide/page counters like '12 / 56' and 'Slide 3'.
    """
    lines = [ln.strip() for ln in text.split("\n")]
    # frequency of short lines
    short = [ln for ln in lines if 3 <= len(ln) <= 80]
    freq = Counter(short)
    banned = set([ln for (ln, c) in freq.items() if c >= repeat_threshold])

    out = []
    for ln in lines:
        if not ln:
            out.append("")
            continue
        if ln in banned:
            continue
        # page/slide counters
        if re.fullmatch(r"\d+\s*/\s*\d+", ln):            # "12 / 56"
            continue
        if re.fullmatch(r"(page|slide)\s*\d+\b.*", ln, flags=re.I):
            continue
        if re.fullmatch(r"week\s+\d+.*", ln, flags=re.I):
            continue
        out.append(ln)
    return "\n".join(out)

def merge_wrapped_lines(text: str) -> str:
    """
    Merge lines that were wrapped by the extractor:
    - keep bullet/numbered lines separate
    - join short lines that are likely mid-sentence
    """
    lines = text.split("\n")
    out = []
    buffer = ""
    bullet_re = re.compile(r"^\s*([•\-–]\s+|\d+[\.)]\s+|[a-z]\)\s+)")
    for raw in lines:
        ln = raw.strip()
        if not ln:
            if buffer:
                out.append(buffer.strip())
                buffer = ""
            out.append("")
            continue
        # new bullet => flush buffer
        if bullet_re.match(ln):
            if buffer:
                out.append(buffer.strip())
                buffer = ""
            out.append(ln)
            continue
        # join heuristics
        if not buffer:
            buffer = ln
        else:
            # if previous ends with sentence end, start new; else join with space
            if re.search(r"[.!?:]$", buffer):
                out.append(buffer)
                buffer = ln
            else:
                # avoid double spaces when hyphenated line breaks
                if buffer.endswith("-"):
                    buffer = buffer[:-1] + ln
                else:
                    buffer += " " + ln
    if buffer:
        out.append(buffer)
    return "\n".join(out)

def compact_whitespace(text: str) -> str:
    # collapse 3+ blanks to single blank
    text = re.sub(r"\n{3,}", "\n\n", text)
    # trim trailing spaces per line
    text = "\n".join(ln.rstrip() for ln in text.splitlines())
    return text.strip()

def compact_text_for_llm(text: str) -> str:
    text = sanitize(text)
    text = drop_headers_footers(text, repeat_threshold=3)
    text = merge_wrapped_lines(text)
    text = compact_whitespace(text)
    return text

def main():
    if len(sys.argv) < 2:
        print("Usage: python extraction.py <file.pdf or file.pptx>")
        sys.exit(1)

    file_path = sys.argv[1]

    if not os.path.isfile(file_path):
        print(f"Error: file not found → {file_path}")
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        text = extract_text(file_path)
    elif ext == ".pptx":
        prs = Presentation(file_path)
        slide_texts = []

        for i, slide in enumerate(prs.slides, start=1):
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content.append(shape.text)
            if content:
                slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(content))

        text = "\n\n".join(slide_texts)
    else:
        print("Error: only .pdf or .pptx files are supported.")
        sys.exit(1)

    text = compact_text_for_llm(text)
    safe_text = text.encode("utf-8", errors="ignore").decode("utf-8")

    # --- Save to output.txt in same folder ---
    output_path = os.path.join(os.path.dirname(file_path), "output.txt")
    with open(output_path, "w", encoding="utf-8", errors="strict") as f:
        f.write(safe_text)

if __name__ == "__main__":
    main()